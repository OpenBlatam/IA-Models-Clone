"""
Advanced Document Generator for Business Agents
=============================================

Comprehensive document generation system with templates, AI assistance, and multi-format support.
"""

import asyncio
import logging
import json
import hashlib
from typing import Dict, Any, List, Optional, Union, Tuple
from datetime import datetime, timedelta
from uuid import UUID, uuid4
from dataclasses import dataclass, field
from enum import Enum
import jinja2
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
import pandas as pd
from docx import Document
from docx.shared import Inches
from pptx import Presentation
from pptx.util import Inches as PptxInches
import matplotlib.pyplot as plt
import seaborn as sns
from sqlalchemy.ext.asyncio import AsyncSession
import redis

from .schemas import (
    BusinessAgent, AgentRequest, AgentResponse, AgentAnalytics,
    AgentWorkflow, AgentCollaboration, AgentSettings,
    ErrorResponse
)
from .exceptions import (
    DocumentGenerationError, TemplateNotFoundError, DocumentValidationError,
    DocumentExportError, DocumentSystemError,
    create_agent_error, log_agent_error, handle_agent_error, get_error_response
)
from .config import get_settings

logger = logging.getLogger(__name__)


class DocumentType(Enum):
    """Document type enumeration"""
    REPORT = "report"
    PROPOSAL = "proposal"
    MANUAL = "manual"
    PRESENTATION = "presentation"
    SPREADSHEET = "spreadsheet"
    DASHBOARD = "dashboard"
    CONTRACT = "contract"
    PROPOSAL = "proposal"
    WHITEPAPER = "whitepaper"
    CASE_STUDY = "case_study"
    USER_GUIDE = "user_guide"
    API_DOCUMENTATION = "api_documentation"
    TECHNICAL_SPEC = "technical_spec"
    BUSINESS_PLAN = "business_plan"
    MARKETING_MATERIAL = "marketing_material"


class DocumentFormat(Enum):
    """Document format enumeration"""
    PDF = "pdf"
    DOCX = "docx"
    PPTX = "pptx"
    XLSX = "xlsx"
    HTML = "html"
    MARKDOWN = "markdown"
    JSON = "json"
    XML = "xml"
    CSV = "csv"
    TXT = "txt"


class TemplateType(Enum):
    """Template type enumeration"""
    STATIC = "static"
    DYNAMIC = "dynamic"
    AI_GENERATED = "ai_generated"
    CUSTOM = "custom"


@dataclass
class DocumentTemplate:
    """Document template definition"""
    template_id: str
    name: str
    description: str
    document_type: DocumentType
    template_type: TemplateType
    format: DocumentFormat
    content: str
    variables: Dict[str, Any] = field(default_factory=dict)
    metadata: Dict[str, Any] = field(default_factory=dict)
    created_by: str = ""
    created_at: datetime = field(default_factory=datetime.utcnow)
    updated_at: datetime = field(default_factory=datetime.utcnow)


@dataclass
class DocumentRequest:
    """Document generation request"""
    request_id: str
    template_id: str
    document_type: DocumentType
    format: DocumentFormat
    data: Dict[str, Any] = field(default_factory=dict)
    variables: Dict[str, Any] = field(default_factory=dict)
    options: Dict[str, Any] = field(default_factory=dict)
    user_id: str = ""
    timestamp: datetime = field(default_factory=datetime.utcnow)


@dataclass
class GeneratedDocument:
    """Generated document definition"""
    document_id: str
    request_id: str
    name: str
    document_type: DocumentType
    format: DocumentFormat
    content: bytes
    metadata: Dict[str, Any] = field(default_factory=dict)
    file_size: int = 0
    generated_at: datetime = field(default_factory=datetime.utcnow)


class DocumentGenerator:
    """Advanced document generator for business agents"""
    
    def __init__(self, db_session: AsyncSession, redis_client: redis.Redis):
        self.db = db_session
        self.redis = redis_client
        self.settings = get_settings()
        self.templates: Dict[str, DocumentTemplate] = {}
        self.jinja_env = jinja2.Environment(loader=jinja2.DictLoader({}))
        self._initialize_templates()
    
    def _initialize_templates(self):
        """Initialize default document templates"""
        try:
            # Agent Report Template
            agent_report_template = DocumentTemplate(
                template_id="agent_report_template",
                name="Agent Performance Report",
                description="Comprehensive agent performance report template",
                document_type=DocumentType.REPORT,
                template_type=TemplateType.STATIC,
                format=DocumentFormat.PDF,
                content=self._get_agent_report_template(),
                variables={
                    "agent_name": "string",
                    "report_period": "string",
                    "performance_metrics": "object",
                    "analytics_data": "object"
                }
            )
            self.templates[agent_report_template.template_id] = agent_report_template
            
            # Workflow Documentation Template
            workflow_doc_template = DocumentTemplate(
                template_id="workflow_doc_template",
                name="Workflow Documentation",
                description="Workflow documentation template",
                document_type=DocumentType.MANUAL,
                template_type=TemplateType.STATIC,
                format=DocumentFormat.DOCX,
                content=self._get_workflow_doc_template(),
                variables={
                    "workflow_name": "string",
                    "workflow_description": "string",
                    "workflow_steps": "array",
                    "workflow_diagram": "string"
                }
            )
            self.templates[workflow_doc_template.template_id] = workflow_doc_template
            
            # Analytics Dashboard Template
            dashboard_template = DocumentTemplate(
                template_id="dashboard_template",
                name="Analytics Dashboard",
                description="Analytics dashboard template",
                document_type=DocumentType.DASHBOARD,
                template_type=TemplateType.DYNAMIC,
                format=DocumentFormat.HTML,
                content=self._get_dashboard_template(),
                variables={
                    "dashboard_title": "string",
                    "metrics_data": "object",
                    "charts_data": "array",
                    "time_period": "string"
                }
            )
            self.templates[dashboard_template.template_id] = dashboard_template
            
            logger.info("Document templates initialized successfully")
            
        except Exception as e:
            logger.error(f"Failed to initialize document templates: {e}")
            raise DocumentGenerationError(
                "template_initialization_failed",
                "Failed to initialize document templates",
                {"error": str(e)}
            )
    
    async def generate_document(
        self,
        template_id: str,
        data: Dict[str, Any],
        format: DocumentFormat = DocumentFormat.PDF,
        options: Dict[str, Any] = None,
        user_id: str = None
    ) -> GeneratedDocument:
        """Generate document from template"""
        try:
            # Get template
            template = await self._get_template(template_id)
            if not template:
                raise TemplateNotFoundError(
                    "template_not_found",
                    f"Template {template_id} not found",
                    {"template_id": template_id}
                )
            
            # Create request
            request = DocumentRequest(
                request_id=str(uuid4()),
                template_id=template_id,
                document_type=template.document_type,
                format=format,
                data=data,
                options=options or {},
                user_id=user_id or ""
            )
            
            # Generate document based on format
            if format == DocumentFormat.PDF:
                content = await self._generate_pdf(template, data, options)
            elif format == DocumentFormat.DOCX:
                content = await self._generate_docx(template, data, options)
            elif format == DocumentFormat.PPTX:
                content = await self._generate_pptx(template, data, options)
            elif format == DocumentFormat.XLSX:
                content = await self._generate_xlsx(template, data, options)
            elif format == DocumentFormat.HTML:
                content = await self._generate_html(template, data, options)
            elif format == DocumentFormat.MARKDOWN:
                content = await self._generate_markdown(template, data, options)
            else:
                raise DocumentGenerationError(
                    "unsupported_format",
                    f"Format {format} not supported",
                    {"format": format.value}
                )
            
            # Create generated document
            document = GeneratedDocument(
                document_id=str(uuid4()),
                request_id=request.request_id,
                name=f"{template.name}_{datetime.utcnow().strftime('%Y%m%d_%H%M%S')}",
                document_type=template.document_type,
                format=format,
                content=content,
                file_size=len(content),
                metadata={
                    "template_id": template_id,
                    "template_name": template.name,
                    "generated_by": user_id,
                    "generation_time": datetime.utcnow().isoformat()
                }
            )
            
            # Store document
            await self._store_document(document)
            
            logger.info(f"Document generated successfully: {document.document_id}")
            
            return document
            
        except Exception as e:
            error = handle_agent_error(e, template_id=template_id, user_id=user_id)
            log_agent_error(error)
            raise error
    
    async def generate_agent_report(
        self,
        agent_id: str,
        report_period: str = "30d",
        format: DocumentFormat = DocumentFormat.PDF,
        user_id: str = None
    ) -> GeneratedDocument:
        """Generate agent performance report"""
        try:
            # Get agent data
            agent_data = await self._get_agent_data(agent_id)
            if not agent_data:
                raise DocumentGenerationError(
                    "agent_not_found",
                    f"Agent {agent_id} not found",
                    {"agent_id": agent_id}
                )
            
            # Get analytics data
            analytics_data = await self._get_agent_analytics(agent_id, report_period)
            
            # Prepare data for template
            data = {
                "agent_name": agent_data.get("name", "Unknown Agent"),
                "agent_type": agent_data.get("agent_type", "Unknown"),
                "report_period": report_period,
                "generated_at": datetime.utcnow().isoformat(),
                "performance_metrics": {
                    "execution_count": analytics_data.get("execution_count", 0),
                    "success_rate": analytics_data.get("success_rate", 0.0),
                    "average_response_time": analytics_data.get("average_response_time", 0.0),
                    "error_rate": analytics_data.get("error_rate", 0.0)
                },
                "analytics_data": analytics_data,
                "agent_configuration": agent_data.get("configuration", {}),
                "trends": analytics_data.get("trends", {})
            }
            
            # Generate document
            document = await self.generate_document(
                template_id="agent_report_template",
                data=data,
                format=format,
                user_id=user_id
            )
            
            logger.info(f"Agent report generated for agent {agent_id}")
            
            return document
            
        except Exception as e:
            error = handle_agent_error(e, agent_id=agent_id, user_id=user_id)
            log_agent_error(error)
            raise error
    
    async def generate_workflow_documentation(
        self,
        workflow_id: str,
        format: DocumentFormat = DocumentFormat.DOCX,
        user_id: str = None
    ) -> GeneratedDocument:
        """Generate workflow documentation"""
        try:
            # Get workflow data
            workflow_data = await self._get_workflow_data(workflow_id)
            if not workflow_data:
                raise DocumentGenerationError(
                    "workflow_not_found",
                    f"Workflow {workflow_id} not found",
                    {"workflow_id": workflow_id}
                )
            
            # Prepare data for template
            data = {
                "workflow_name": workflow_data.get("name", "Unknown Workflow"),
                "workflow_description": workflow_data.get("description", ""),
                "workflow_type": workflow_data.get("workflow_type", "Unknown"),
                "workflow_steps": workflow_data.get("nodes", []),
                "workflow_connections": workflow_data.get("connections", []),
                "workflow_variables": workflow_data.get("variables", {}),
                "created_at": workflow_data.get("created_at", ""),
                "updated_at": workflow_data.get("updated_at", ""),
                "created_by": workflow_data.get("created_by", "")
            }
            
            # Generate document
            document = await self.generate_document(
                template_id="workflow_doc_template",
                data=data,
                format=format,
                user_id=user_id
            )
            
            logger.info(f"Workflow documentation generated for workflow {workflow_id}")
            
            return document
            
        except Exception as e:
            error = handle_agent_error(e, workflow_id=workflow_id, user_id=user_id)
            log_agent_error(error)
            raise error
    
    async def generate_analytics_dashboard(
        self,
        dashboard_type: str,
        time_period: str = "30d",
        format: DocumentFormat = DocumentFormat.HTML,
        user_id: str = None
    ) -> GeneratedDocument:
        """Generate analytics dashboard"""
        try:
            # Get analytics data
            analytics_data = await self._get_analytics_data(dashboard_type, time_period)
            
            # Prepare data for template
            data = {
                "dashboard_title": f"{dashboard_type.title()} Analytics Dashboard",
                "time_period": time_period,
                "generated_at": datetime.utcnow().isoformat(),
                "metrics_data": analytics_data.get("metrics", {}),
                "charts_data": analytics_data.get("charts", []),
                "insights": analytics_data.get("insights", []),
                "recommendations": analytics_data.get("recommendations", [])
            }
            
            # Generate document
            document = await self.generate_document(
                template_id="dashboard_template",
                data=data,
                format=format,
                user_id=user_id
            )
            
            logger.info(f"Analytics dashboard generated for {dashboard_type}")
            
            return document
            
        except Exception as e:
            error = handle_agent_error(e, dashboard_type=dashboard_type, user_id=user_id)
            log_agent_error(error)
            raise error
    
    async def create_custom_template(
        self,
        name: str,
        description: str,
        document_type: DocumentType,
        format: DocumentFormat,
        content: str,
        variables: Dict[str, Any] = None,
        user_id: str = None
    ) -> DocumentTemplate:
        """Create custom document template"""
        try:
            # Create template
            template = DocumentTemplate(
                template_id=str(uuid4()),
                name=name,
                description=description,
                document_type=document_type,
                template_type=TemplateType.CUSTOM,
                format=format,
                content=content,
                variables=variables or {},
                created_by=user_id or ""
            )
            
            # Store template
            await self._store_template(template)
            
            # Add to local cache
            self.templates[template.template_id] = template
            
            logger.info(f"Custom template created: {template.template_id}")
            
            return template
            
        except Exception as e:
            error = handle_agent_error(e, user_id=user_id)
            log_agent_error(error)
            raise error
    
    async def list_templates(
        self,
        document_type: DocumentType = None,
        template_type: TemplateType = None
    ) -> List[DocumentTemplate]:
        """List available templates"""
        try:
            templates = list(self.templates.values())
            
            # Filter by document type
            if document_type:
                templates = [t for t in templates if t.document_type == document_type]
            
            # Filter by template type
            if template_type:
                templates = [t for t in templates if t.template_type == template_type]
            
            return templates
            
        except Exception as e:
            logger.error(f"Failed to list templates: {e}")
            raise DocumentGenerationError(
                "template_listing_failed",
                "Failed to list templates",
                {"error": str(e)}
            )
    
    async def get_template(self, template_id: str) -> Optional[DocumentTemplate]:
        """Get template by ID"""
        try:
            return await self._get_template(template_id)
        except Exception as e:
            logger.error(f"Failed to get template: {e}")
            raise DocumentGenerationError(
                "template_retrieval_failed",
                "Failed to get template",
                {"error": str(e)}
            )
    
    # Document generation methods
    async def _generate_pdf(self, template: DocumentTemplate, data: Dict[str, Any], options: Dict[str, Any]) -> bytes:
        """Generate PDF document"""
        try:
            from io import BytesIO
            
            buffer = BytesIO()
            doc = SimpleDocTemplate(buffer, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            # Title
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=18,
                spaceAfter=30,
                alignment=1  # Center alignment
            )
            story.append(Paragraph(template.name, title_style))
            story.append(Spacer(1, 12))
            
            # Process template content
            content = self._process_template(template.content, data)
            
            # Add content paragraphs
            for line in content.split('\n'):
                if line.strip():
                    story.append(Paragraph(line, styles['Normal']))
                    story.append(Spacer(1, 6))
            
            # Add performance metrics table if available
            if 'performance_metrics' in data:
                metrics_data = data['performance_metrics']
                table_data = [
                    ['Metric', 'Value'],
                    ['Execution Count', str(metrics_data.get('execution_count', 0))],
                    ['Success Rate', f"{metrics_data.get('success_rate', 0.0):.2%}"],
                    ['Average Response Time', f"{metrics_data.get('average_response_time', 0.0):.2f}s"],
                    ['Error Rate', f"{metrics_data.get('error_rate', 0.0):.2%}"]
                ]
                
                table = Table(table_data)
                table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 14),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black)
                ]))
                
                story.append(Spacer(1, 20))
                story.append(table)
            
            # Build PDF
            doc.build(story)
            buffer.seek(0)
            
            return buffer.getvalue()
            
        except Exception as e:
            logger.error(f"Failed to generate PDF: {e}")
            raise DocumentExportError(
                "pdf_generation_failed",
                "Failed to generate PDF document",
                {"error": str(e)}
            )
    
    async def _generate_docx(self, template: DocumentTemplate, data: Dict[str, Any], options: Dict[str, Any]) -> bytes:
        """Generate DOCX document"""
        try:
            from io import BytesIO
            
            doc = Document()
            
            # Add title
            doc.add_heading(template.name, 0)
            
            # Process template content
            content = self._process_template(template.content, data)
            
            # Add content
            for line in content.split('\n'):
                if line.strip():
                    doc.add_paragraph(line)
            
            # Add performance metrics table if available
            if 'performance_metrics' in data:
                doc.add_heading('Performance Metrics', level=1)
                metrics_data = data['performance_metrics']
                
                table = doc.add_table(rows=1, cols=2)
                table.style = 'Table Grid'
                hdr_cells = table.rows[0].cells
                hdr_cells[0].text = 'Metric'
                hdr_cells[1].text = 'Value'
                
                metrics = [
                    ('Execution Count', str(metrics_data.get('execution_count', 0))),
                    ('Success Rate', f"{metrics_data.get('success_rate', 0.0):.2%}"),
                    ('Average Response Time', f"{metrics_data.get('average_response_time', 0.0):.2f}s"),
                    ('Error Rate', f"{metrics_data.get('error_rate', 0.0):.2%}")
                ]
                
                for metric, value in metrics:
                    row_cells = table.add_row().cells
                    row_cells[0].text = metric
                    row_cells[1].text = value
            
            # Save to buffer
            buffer = BytesIO()
            doc.save(buffer)
            buffer.seek(0)
            
            return buffer.getvalue()
            
        except Exception as e:
            logger.error(f"Failed to generate DOCX: {e}")
            raise DocumentExportError(
                "docx_generation_failed",
                "Failed to generate DOCX document",
                {"error": str(e)}
            )
    
    async def _generate_pptx(self, template: DocumentTemplate, data: Dict[str, Any], options: Dict[str, Any]) -> bytes:
        """Generate PPTX presentation"""
        try:
            from io import BytesIO
            
            prs = Presentation()
            
            # Title slide
            slide_layout = prs.slide_layouts[0]  # Title slide layout
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = template.name
            subtitle.text = f"Generated on {datetime.utcnow().strftime('%Y-%m-%d %H:%M:%S')}"
            
            # Content slide
            slide_layout = prs.slide_layouts[1]  # Content slide layout
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = "Overview"
            
            # Process template content
            template_content = self._process_template(template.content, data)
            content.text = template_content[:500] + "..." if len(template_content) > 500 else template_content
            
            # Performance metrics slide
            if 'performance_metrics' in data:
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                content = slide.placeholders[1]
                
                title.text = "Performance Metrics"
                
                metrics_data = data['performance_metrics']
                metrics_text = f"""
                Execution Count: {metrics_data.get('execution_count', 0)}
                Success Rate: {metrics_data.get('success_rate', 0.0):.2%}
                Average Response Time: {metrics_data.get('average_response_time', 0.0):.2f}s
                Error Rate: {metrics_data.get('error_rate', 0.0):.2%}
                """
                content.text = metrics_text
            
            # Save to buffer
            buffer = BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            
            return buffer.getvalue()
            
        except Exception as e:
            logger.error(f"Failed to generate PPTX: {e}")
            raise DocumentExportError(
                "pptx_generation_failed",
                "Failed to generate PPTX presentation",
                {"error": str(e)}
            )
    
    async def _generate_xlsx(self, template: DocumentTemplate, data: Dict[str, Any], options: Dict[str, Any]) -> bytes:
        """Generate XLSX spreadsheet"""
        try:
            from io import BytesIO
            
            # Create DataFrame with analytics data
            if 'analytics_data' in data:
                analytics = data['analytics_data']
                df = pd.DataFrame(analytics.get('daily_metrics', []))
            else:
                # Create sample data
                df = pd.DataFrame({
                    'Date': pd.date_range(start='2024-01-01', periods=30, freq='D'),
                    'Executions': range(100, 130),
                    'Success Rate': [0.95 + (i * 0.001) for i in range(30)],
                    'Response Time': [2.0 + (i * 0.1) for i in range(30)]
                })
            
            # Save to buffer
            buffer = BytesIO()
            with pd.ExcelWriter(buffer, engine='openpyxl') as writer:
                df.to_excel(writer, sheet_name='Analytics', index=False)
                
                # Add summary sheet
                summary_data = {
                    'Metric': ['Total Executions', 'Average Success Rate', 'Average Response Time'],
                    'Value': [
                        df['Executions'].sum() if 'Executions' in df.columns else 0,
                        df['Success Rate'].mean() if 'Success Rate' in df.columns else 0,
                        df['Response Time'].mean() if 'Response Time' in df.columns else 0
                    ]
                }
                summary_df = pd.DataFrame(summary_data)
                summary_df.to_excel(writer, sheet_name='Summary', index=False)
            
            buffer.seek(0)
            return buffer.getvalue()
            
        except Exception as e:
            logger.error(f"Failed to generate XLSX: {e}")
            raise DocumentExportError(
                "xlsx_generation_failed",
                "Failed to generate XLSX spreadsheet",
                {"error": str(e)}
            )
    
    async def _generate_html(self, template: DocumentTemplate, data: Dict[str, Any], options: Dict[str, Any]) -> bytes:
        """Generate HTML document"""
        try:
            # Process template content
            content = self._process_template(template.content, data)
            
            # Create HTML document
            html_content = f"""
            <!DOCTYPE html>
            <html lang="en">
            <head>
                <meta charset="UTF-8">
                <meta name="viewport" content="width=device-width, initial-scale=1.0">
                <title>{template.name}</title>
                <style>
                    body {{ font-family: Arial, sans-serif; margin: 40px; }}
                    h1 {{ color: #333; border-bottom: 2px solid #333; }}
                    h2 {{ color: #666; }}
                    table {{ border-collapse: collapse; width: 100%; }}
                    th, td {{ border: 1px solid #ddd; padding: 8px; text-align: left; }}
                    th {{ background-color: #f2f2f2; }}
                    .metric {{ background-color: #f9f9f9; padding: 10px; margin: 10px 0; }}
                </style>
            </head>
            <body>
                <h1>{template.name}</h1>
                <p>Generated on: {datetime.utcnow().strftime('%Y-%m-%d %H:%M:%S')}</p>
                
                <div class="content">
                    {content.replace(chr(10), '<br>')}
                </div>
                
                {self._generate_metrics_html(data) if 'performance_metrics' in data else ''}
                
                {self._generate_charts_html(data) if 'charts_data' in data else ''}
            </body>
            </html>
            """
            
            return html_content.encode('utf-8')
            
        except Exception as e:
            logger.error(f"Failed to generate HTML: {e}")
            raise DocumentExportError(
                "html_generation_failed",
                "Failed to generate HTML document",
                {"error": str(e)}
            )
    
    async def _generate_markdown(self, template: DocumentTemplate, data: Dict[str, Any], options: Dict[str, Any]) -> bytes:
        """Generate Markdown document"""
        try:
            # Process template content
            content = self._process_template(template.content, data)
            
            # Create Markdown document
            markdown_content = f"""# {template.name}

**Generated on:** {datetime.utcnow().strftime('%Y-%m-%d %H:%M:%S')}

---

{content}

{self._generate_metrics_markdown(data) if 'performance_metrics' in data else ''}

{self._generate_insights_markdown(data) if 'insights' in data else ''}
"""
            
            return markdown_content.encode('utf-8')
            
        except Exception as e:
            logger.error(f"Failed to generate Markdown: {e}")
            raise DocumentExportError(
                "markdown_generation_failed",
                "Failed to generate Markdown document",
                {"error": str(e)}
            )
    
    # Helper methods
    def _process_template(self, template_content: str, data: Dict[str, Any]) -> str:
        """Process template with data"""
        try:
            # Use Jinja2 for template processing
            template = self.jinja_env.from_string(template_content)
            return template.render(**data)
        except Exception as e:
            logger.error(f"Failed to process template: {e}")
            return template_content
    
    def _generate_metrics_html(self, data: Dict[str, Any]) -> str:
        """Generate metrics HTML"""
        metrics = data.get('performance_metrics', {})
        return f"""
        <h2>Performance Metrics</h2>
        <table>
            <tr><th>Metric</th><th>Value</th></tr>
            <tr><td>Execution Count</td><td>{metrics.get('execution_count', 0)}</td></tr>
            <tr><td>Success Rate</td><td>{metrics.get('success_rate', 0.0):.2%}</td></tr>
            <tr><td>Average Response Time</td><td>{metrics.get('average_response_time', 0.0):.2f}s</td></tr>
            <tr><td>Error Rate</td><td>{metrics.get('error_rate', 0.0):.2%}</td></tr>
        </table>
        """
    
    def _generate_metrics_markdown(self, data: Dict[str, Any]) -> str:
        """Generate metrics Markdown"""
        metrics = data.get('performance_metrics', {})
        return f"""
## Performance Metrics

| Metric | Value |
|--------|-------|
| Execution Count | {metrics.get('execution_count', 0)} |
| Success Rate | {metrics.get('success_rate', 0.0):.2%} |
| Average Response Time | {metrics.get('average_response_time', 0.0):.2f}s |
| Error Rate | {metrics.get('error_rate', 0.0):.2%} |
"""
    
    def _generate_charts_html(self, data: Dict[str, Any]) -> str:
        """Generate charts HTML"""
        charts = data.get('charts_data', [])
        if not charts:
            return ""
        
        html = "<h2>Charts</h2>"
        for chart in charts:
            html += f"<div class='chart'>{chart.get('title', 'Chart')}</div>"
        
        return html
    
    def _generate_insights_markdown(self, data: Dict[str, Any]) -> str:
        """Generate insights Markdown"""
        insights = data.get('insights', [])
        if not insights:
            return ""
        
        markdown = "## Insights\n\n"
        for insight in insights:
            markdown += f"- {insight}\n"
        
        return markdown
    
    # Template content methods
    def _get_agent_report_template(self) -> str:
        """Get agent report template"""
        return """
# Agent Performance Report

## Agent Information
- **Name:** {{ agent_name }}
- **Type:** {{ agent_type }}
- **Report Period:** {{ report_period }}
- **Generated:** {{ generated_at }}

## Performance Summary
The agent has shown the following performance during the reporting period:

- **Total Executions:** {{ performance_metrics.execution_count }}
- **Success Rate:** {{ performance_metrics.success_rate }}%
- **Average Response Time:** {{ performance_metrics.average_response_time }}s
- **Error Rate:** {{ performance_metrics.error_rate }}%

## Detailed Analysis
{% if analytics_data.trends %}
### Performance Trends
The agent's performance has shown the following trends:
- Execution volume has {{ analytics_data.trends.execution_trend }}
- Success rate has {{ analytics_data.trends.success_trend }}
- Response time has {{ analytics_data.trends.response_trend }}
{% endif %}

## Recommendations
Based on the analysis, the following recommendations are made:
1. Monitor performance metrics regularly
2. Optimize configuration for better performance
3. Review error patterns and implement fixes
"""
    
    def _get_workflow_doc_template(self) -> str:
        """Get workflow documentation template"""
        return """
# Workflow Documentation

## Workflow Information
- **Name:** {{ workflow_name }}
- **Type:** {{ workflow_type }}
- **Description:** {{ workflow_description }}
- **Created:** {{ created_at }}
- **Last Updated:** {{ updated_at }}

## Workflow Steps
{% for step in workflow_steps %}
### Step {{ loop.index }}: {{ step.name }}
- **Type:** {{ step.node_type }}
- **Description:** {{ step.description }}
- **Configuration:** {{ step.properties }}
{% endfor %}

## Workflow Connections
{% for connection in workflow_connections %}
- **From:** {{ connection.source_node }}
- **To:** {{ connection.target_node }}
- **Type:** {{ connection.connection_type }}
{% if connection.condition %}
- **Condition:** {{ connection.condition }}
{% endif %}
{% endfor %}

## Variables
{% for key, value in workflow_variables.items() %}
- **{{ key }}:** {{ value }}
{% endfor %}
"""
    
    def _get_dashboard_template(self) -> str:
        """Get dashboard template"""
        return """
<!DOCTYPE html>
<html>
<head>
    <title>{{ dashboard_title }}</title>
    <script src="https://cdn.jsdelivr.net/npm/chart.js"></script>
    <style>
        body { font-family: Arial, sans-serif; margin: 20px; }
        .metric-card { border: 1px solid #ddd; padding: 15px; margin: 10px; display: inline-block; }
        .chart-container { width: 400px; height: 300px; margin: 20px; }
    </style>
</head>
<body>
    <h1>{{ dashboard_title }}</h1>
    <p>Period: {{ time_period }} | Generated: {{ generated_at }}</p>
    
    <h2>Key Metrics</h2>
    {% for key, value in metrics_data.items() %}
    <div class="metric-card">
        <h3>{{ key }}</h3>
        <p>{{ value }}</p>
    </div>
    {% endfor %}
    
    <h2>Charts</h2>
    {% for chart in charts_data %}
    <div class="chart-container">
        <canvas id="chart{{ loop.index }}"></canvas>
    </div>
    {% endfor %}
    
    {% if insights %}
    <h2>Insights</h2>
    <ul>
    {% for insight in insights %}
        <li>{{ insight }}</li>
    {% endfor %}
    </ul>
    {% endif %}
    
    {% if recommendations %}
    <h2>Recommendations</h2>
    <ul>
    {% for recommendation in recommendations %}
        <li>{{ recommendation }}</li>
    {% endfor %}
    </ul>
    {% endif %}
</body>
</html>
"""
    
    # Data retrieval methods
    async def _get_agent_data(self, agent_id: str) -> Optional[Dict[str, Any]]:
        """Get agent data"""
        try:
            # This would integrate with actual database
            # For now, return mock data
            return {
                "id": agent_id,
                "name": f"Agent {agent_id}",
                "agent_type": "sales",
                "configuration": {"param1": "value1"},
                "created_at": datetime.utcnow().isoformat()
            }
        except Exception as e:
            logger.error(f"Failed to get agent data: {e}")
            return None
    
    async def _get_workflow_data(self, workflow_id: str) -> Optional[Dict[str, Any]]:
        """Get workflow data"""
        try:
            # This would integrate with actual database
            # For now, return mock data
            return {
                "id": workflow_id,
                "name": f"Workflow {workflow_id}",
                "workflow_type": "automation",
                "description": "Sample workflow",
                "nodes": [],
                "connections": [],
                "variables": {},
                "created_at": datetime.utcnow().isoformat(),
                "updated_at": datetime.utcnow().isoformat()
            }
        except Exception as e:
            logger.error(f"Failed to get workflow data: {e}")
            return None
    
    async def _get_agent_analytics(self, agent_id: str, period: str) -> Dict[str, Any]:
        """Get agent analytics data"""
        try:
            # This would integrate with actual analytics service
            # For now, return mock data
            return {
                "execution_count": 150,
                "success_rate": 0.95,
                "average_response_time": 2.5,
                "error_rate": 0.05,
                "trends": {
                    "execution_trend": "increased by 10%",
                    "success_trend": "remained stable",
                    "response_trend": "improved by 5%"
                }
            }
        except Exception as e:
            logger.error(f"Failed to get agent analytics: {e}")
            return {}
    
    async def _get_analytics_data(self, dashboard_type: str, period: str) -> Dict[str, Any]:
        """Get analytics data"""
        try:
            # This would integrate with actual analytics service
            # For now, return mock data
            return {
                "metrics": {
                    "Total Agents": 25,
                    "Active Workflows": 15,
                    "Success Rate": "95%",
                    "Average Response Time": "2.5s"
                },
                "charts": [
                    {"title": "Performance Trends", "type": "line"},
                    {"title": "Agent Distribution", "type": "pie"}
                ],
                "insights": [
                    "System performance is optimal",
                    "Agent utilization is high"
                ],
                "recommendations": [
                    "Consider scaling infrastructure",
                    "Optimize workflow configurations"
                ]
            }
        except Exception as e:
            logger.error(f"Failed to get analytics data: {e}")
            return {}
    
    async def _get_template(self, template_id: str) -> Optional[DocumentTemplate]:
        """Get template by ID"""
        try:
            # Try cache first
            cached_template = await self._get_cached_template(template_id)
            if cached_template:
                return cached_template
            
            # Get from local templates
            return self.templates.get(template_id)
            
        except Exception as e:
            logger.error(f"Failed to get template: {e}")
            return None
    
    async def _store_template(self, template: DocumentTemplate) -> None:
        """Store template"""
        try:
            # This would integrate with actual database
            # For now, just log
            logger.info(f"Storing template: {template.template_id}")
            
        except Exception as e:
            logger.error(f"Failed to store template: {e}")
            raise
    
    async def _store_document(self, document: GeneratedDocument) -> None:
        """Store generated document"""
        try:
            # This would integrate with actual storage
            # For now, just log
            logger.info(f"Storing document: {document.document_id}")
            
        except Exception as e:
            logger.error(f"Failed to store document: {e}")
            raise
    
    async def _get_cached_template(self, template_id: str) -> Optional[DocumentTemplate]:
        """Get cached template"""
        try:
            cache_key = f"template:{template_id}"
            cached_data = await self.redis.get(cache_key)
            if cached_data:
                template_data = json.loads(cached_data)
                return DocumentTemplate(**template_data)
            
        except Exception as e:
            logger.error(f"Failed to get cached template: {e}")
        
        return None