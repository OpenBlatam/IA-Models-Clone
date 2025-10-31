"""
Document Formatters
==================

Format-specific document generators and formatters.
"""

import logging
from abc import ABC, abstractmethod
from typing import Dict, Any, Optional, List
from pathlib import Path
from datetime import datetime

from .types import DocumentTemplate, GeneratedDocument, DocumentFormat

logger = logging.getLogger(__name__)

class BaseFormatter(ABC):
    """Base class for document formatters."""
    
    def __init__(self, output_dir: str = "generated_documents"):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(exist_ok=True)
    
    @abstractmethod
    async def format_document(
        self, 
        template: DocumentTemplate, 
        variables: Dict[str, Any],
        filename: str
    ) -> GeneratedDocument:
        """Format a document using the template and variables."""
        pass
    
    def _get_file_path(self, filename: str, extension: str) -> Path:
        """Get the full file path for a document."""
        return self.output_dir / f"{filename}.{extension}"

class MarkdownFormatter(BaseFormatter):
    """Markdown document formatter."""
    
    async def format_document(
        self, 
        template: DocumentTemplate, 
        variables: Dict[str, Any],
        filename: str
    ) -> GeneratedDocument:
        """Format a markdown document."""
        try:
            # Simple template variable replacement
            content = template.content
            for key, value in variables.items():
                placeholder = f"{{{{ {key} }}}}"
                content = content.replace(placeholder, str(value))
            
            # Create file
            file_path = self._get_file_path(filename, "md")
            with open(file_path, 'w', encoding='utf-8') as f:
                f.write(content)
            
            return GeneratedDocument(
                document_id=f"md_{filename}_{datetime.now().strftime('%Y%m%d_%H%M%S')}",
                request_id="",  # Will be set by caller
                title=variables.get("title", filename),
                content=content,
                file_path=str(file_path),
                format=DocumentFormat.MARKDOWN,
                size_bytes=len(content.encode('utf-8'))
            )
            
        except Exception as e:
            logger.error(f"Failed to format markdown document: {str(e)}")
            raise

class HTMLFormatter(BaseFormatter):
    """HTML document formatter."""
    
    async def format_document(
        self, 
        template: DocumentTemplate, 
        variables: Dict[str, Any],
        filename: str
    ) -> GeneratedDocument:
        """Format an HTML document."""
        try:
            # Simple template variable replacement
            content = template.content
            for key, value in variables.items():
                placeholder = f"{{{{ {key} }}}}"
                content = content.replace(placeholder, str(value))
            
            # Wrap in basic HTML structure if not already wrapped
            if not content.strip().startswith('<'):
                html_content = f"""<!DOCTYPE html>
<html>
<head>
    <title>{variables.get('title', filename)}</title>
    <meta charset="utf-8">
    <style>
        body {{ font-family: Arial, sans-serif; margin: 40px; }}
        h1, h2, h3 {{ color: #333; }}
        p {{ line-height: 1.6; }}
    </style>
</head>
<body>
{content}
</body>
</html>"""
            else:
                html_content = content
            
            # Create file
            file_path = self._get_file_path(filename, "html")
            with open(file_path, 'w', encoding='utf-8') as f:
                f.write(html_content)
            
            return GeneratedDocument(
                document_id=f"html_{filename}_{datetime.now().strftime('%Y%m%d_%H%M%S')}",
                request_id="",  # Will be set by caller
                title=variables.get("title", filename),
                content=html_content,
                file_path=str(file_path),
                format=DocumentFormat.HTML,
                size_bytes=len(html_content.encode('utf-8'))
            )
            
        except Exception as e:
            logger.error(f"Failed to format HTML document: {str(e)}")
            raise

class JSONFormatter(BaseFormatter):
    """JSON document formatter."""
    
    async def format_document(
        self, 
        template: DocumentTemplate, 
        variables: Dict[str, Any],
        filename: str
    ) -> GeneratedDocument:
        """Format a JSON document."""
        try:
            import json
            
            # Create structured JSON document
            json_content = {
                "title": variables.get("title", filename),
                "template_id": template.template_id,
                "document_type": template.document_type.value,
                "generated_at": datetime.now().isoformat(),
                "content": variables,
                "metadata": template.metadata
            }
            
            # Create file
            file_path = self._get_file_path(filename, "json")
            with open(file_path, 'w', encoding='utf-8') as f:
                json.dump(json_content, f, indent=2, ensure_ascii=False)
            
            content_str = json.dumps(json_content, indent=2, ensure_ascii=False)
            
            return GeneratedDocument(
                document_id=f"json_{filename}_{datetime.now().strftime('%Y%m%d_%H%M%S')}",
                request_id="",  # Will be set by caller
                title=variables.get("title", filename),
                content=content_str,
                file_path=str(file_path),
                format=DocumentFormat.JSON,
                size_bytes=len(content_str.encode('utf-8'))
            )
            
        except Exception as e:
            logger.error(f"Failed to format JSON document: {str(e)}")
            raise

class PDFFormatter(BaseFormatter):
    """PDF document formatter using ReportLab."""
    
    async def format_document(
        self, 
        template: DocumentTemplate, 
        variables: Dict[str, Any],
        filename: str
    ) -> GeneratedDocument:
        """Format a PDF document."""
        try:
            from reportlab.lib.pagesizes import letter, A4
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import inch
            
            # Simple template variable replacement
            content = template.content
            for key, value in variables.items():
                placeholder = f"{{{{ {key} }}}}"
                content = content.replace(placeholder, str(value))
            
            # Create PDF
            file_path = self._get_file_path(filename, "pdf")
            doc = SimpleDocTemplate(str(file_path), pagesize=A4)
            
            # Get styles
            styles = getSampleStyleSheet()
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=18,
                spaceAfter=30,
            )
            
            # Build content
            story = []
            
            # Add title
            title = variables.get("title", filename)
            story.append(Paragraph(title, title_style))
            story.append(Spacer(1, 12))
            
            # Add content (simple text processing)
            paragraphs = content.split('\n\n')
            for para in paragraphs:
                if para.strip():
                    story.append(Paragraph(para.strip(), styles['Normal']))
                    story.append(Spacer(1, 12))
            
            # Build PDF
            doc.build(story)
            
            # Get file size
            size_bytes = file_path.stat().st_size
            
            return GeneratedDocument(
                document_id=f"pdf_{filename}_{datetime.now().strftime('%Y%m%d_%H%M%S')}",
                request_id="",  # Will be set by caller
                title=variables.get("title", filename),
                content=content,
                file_path=str(file_path),
                format=DocumentFormat.PDF,
                size_bytes=size_bytes
            )
            
        except Exception as e:
            logger.error(f"Failed to format PDF document: {str(e)}")
            raise

class DOCXFormatter(BaseFormatter):
    """DOCX document formatter using python-docx."""
    
    async def format_document(
        self, 
        template: DocumentTemplate, 
        variables: Dict[str, Any],
        filename: str
    ) -> GeneratedDocument:
        """Format a DOCX document."""
        try:
            from docx import Document
            from docx.shared import Inches
            
            # Simple template variable replacement
            content = template.content
            for key, value in variables.items():
                placeholder = f"{{{{ {key} }}}}"
                content = content.replace(placeholder, str(value))
            
            # Create DOCX document
            doc = Document()
            
            # Add title
            title = variables.get("title", filename)
            doc.add_heading(title, 0)
            
            # Add content (simple text processing)
            paragraphs = content.split('\n\n')
            for para in paragraphs:
                if para.strip():
                    if para.strip().startswith('#'):
                        # Handle headings
                        level = len(para) - len(para.lstrip('#'))
                        heading_text = para.strip('#').strip()
                        doc.add_heading(heading_text, level)
                    else:
                        doc.add_paragraph(para.strip())
            
            # Save document
            file_path = self._get_file_path(filename, "docx")
            doc.save(str(file_path))
            
            # Get file size
            size_bytes = file_path.stat().st_size
            
            return GeneratedDocument(
                document_id=f"docx_{filename}_{datetime.now().strftime('%Y%m%d_%H%M%S')}",
                request_id="",  # Will be set by caller
                title=variables.get("title", filename),
                content=content,
                file_path=str(file_path),
                format=DocumentFormat.DOCX,
                size_bytes=size_bytes
            )
            
        except Exception as e:
            logger.error(f"Failed to format DOCX document: {str(e)}")
            raise

class PPTXFormatter(BaseFormatter):
    """PPTX presentation formatter using python-pptx."""
    
    async def format_document(
        self, 
        template: DocumentTemplate, 
        variables: Dict[str, Any],
        filename: str
    ) -> GeneratedDocument:
        """Format a PPTX presentation."""
        try:
            from pptx import Presentation
            from pptx.util import Inches
            
            # Simple template variable replacement
            content = template.content
            for key, value in variables.items():
                placeholder = f"{{{{ {key} }}}}"
                content = content.replace(placeholder, str(value))
            
            # Create presentation
            prs = Presentation()
            
            # Add title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = variables.get("title", filename)
            subtitle.text = "Generated by Business Agents System"
            
            # Add content slides (simple text processing)
            paragraphs = content.split('\n\n')
            for para in paragraphs:
                if para.strip() and not para.strip().startswith('#'):
                    # Add content slide
                    bullet_slide_layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(bullet_slide_layout)
                    shapes = slide.shapes
                    
                    title_shape = shapes.title
                    body_shape = shapes.placeholders[1]
                    
                    title_shape.text = "Content"
                    tf = body_shape.text_frame
                    tf.text = para.strip()
            
            # Save presentation
            file_path = self._get_file_path(filename, "pptx")
            prs.save(str(file_path))
            
            # Get file size
            size_bytes = file_path.stat().st_size
            
            return GeneratedDocument(
                document_id=f"pptx_{filename}_{datetime.now().strftime('%Y%m%d_%H%M%S')}",
                request_id="",  # Will be set by caller
                title=variables.get("title", filename),
                content=content,
                file_path=str(file_path),
                format=DocumentFormat.PPTX,
                size_bytes=size_bytes
            )
            
        except Exception as e:
            logger.error(f"Failed to format PPTX document: {str(e)}")
            raise

class FormatterFactory:
    """Factory for creating document formatters."""
    
    _formatters = {
        DocumentFormat.MARKDOWN: MarkdownFormatter,
        DocumentFormat.HTML: HTMLFormatter,
        DocumentFormat.JSON: JSONFormatter,
        DocumentFormat.PDF: PDFFormatter,
        DocumentFormat.DOCX: DOCXFormatter,
        DocumentFormat.PPTX: PPTXFormatter,
    }
    
    @classmethod
    def get_formatter(cls, format: DocumentFormat, output_dir: str = "generated_documents") -> BaseFormatter:
        """Get a formatter for the specified format."""
        formatter_class = cls._formatters.get(format)
        if not formatter_class:
            raise ValueError(f"No formatter available for format: {format}")
        
        return formatter_class(output_dir)
    
    @classmethod
    def get_supported_formats(cls) -> List[DocumentFormat]:
        """Get list of supported formats."""
        return list(cls._formatters.keys())
