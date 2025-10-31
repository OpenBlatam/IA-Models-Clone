"""
Document Upload Processor
Real, working document upload and processing capabilities
"""

import asyncio
import logging
import json
import hashlib
from typing import Dict, List, Any, Optional
from datetime import datetime
import re
import time
import os
import mimetypes
from pathlib import Path

# Real, working libraries for document processing
try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from PIL import Image
    import pytesseract
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

logger = logging.getLogger(__name__)

class DocumentUploadProcessor:
    """Real working document upload and processing"""
    
    def __init__(self):
        self.supported_formats = {
            'pdf': PDF_AVAILABLE,
            'docx': DOCX_AVAILABLE,
            'xlsx': EXCEL_AVAILABLE,
            'pptx': PPTX_AVAILABLE,
            'txt': True,
            'image': OCR_AVAILABLE
        }
        
        # Processing stats
        self.stats = {
            "total_uploads": 0,
            "successful_uploads": 0,
            "failed_uploads": 0,
            "average_processing_time": 0,
            "start_time": time.time()
        }
    
    async def process_uploaded_document(self, file_content: bytes, filename: str, 
                                      file_type: str = None) -> Dict[str, Any]:
        """Process uploaded document and extract text"""
        start_time = time.time()
        
        try:
            logger.info(f"Processing uploaded document: {filename}")
            
            # Generate document ID
            doc_id = self._generate_document_id(file_content, filename)
            
            result = {
                "document_id": doc_id,
                "filename": filename,
                "file_type": file_type or mimetypes.guess_type(filename)[0],
                "file_size": len(file_content),
                "upload_timestamp": datetime.now().isoformat(),
                "processing_successful": False,
                "text_content": "",
                "metadata": {},
                "error": None
            }
            
            # Parse document based on file type
            if filename.lower().endswith('.pdf') and PDF_AVAILABLE:
                result.update(await self._parse_pdf(file_content))
            elif filename.lower().endswith('.docx') and DOCX_AVAILABLE:
                result.update(await self._parse_docx(file_content))
            elif filename.lower().endswith(('.xlsx', '.xls')) and EXCEL_AVAILABLE:
                result.update(await self._parse_excel(file_content))
            elif filename.lower().endswith('.pptx') and PPTX_AVAILABLE:
                result.update(await self._parse_pptx(file_content))
            elif filename.lower().endswith('.txt'):
                result.update(await self._parse_txt(file_content))
            elif self._is_image_file(filename) and OCR_AVAILABLE:
                result.update(await self._parse_image(file_content))
            else:
                result["error"] = f"Unsupported file format: {filename}"
                return result
            
            result["processing_successful"] = True
            
            # Update stats
            processing_time = time.time() - start_time
            self._update_stats(processing_time, True)
            result["processing_time"] = processing_time
            
            return result
            
        except Exception as e:
            processing_time = time.time() - start_time
            self._update_stats(processing_time, False)
            logger.error(f"Error processing uploaded document {filename}: {e}")
            return {
                "document_id": self._generate_document_id(file_content, filename),
                "filename": filename,
                "processing_successful": False,
                "error": str(e),
                "processing_time": processing_time,
                "timestamp": datetime.now().isoformat()
            }
    
    async def _parse_pdf(self, content: bytes) -> Dict[str, Any]:
        """Parse PDF document"""
        try:
            import io
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
            
            text_content = ""
            metadata = {
                "num_pages": len(pdf_reader.pages),
                "title": "",
                "author": "",
                "subject": "",
                "creator": "",
                "producer": "",
                "creation_date": "",
                "modification_date": ""
            }
            
            # Extract metadata
            if pdf_reader.metadata:
                metadata.update({
                    "title": pdf_reader.metadata.get("/Title", ""),
                    "author": pdf_reader.metadata.get("/Author", ""),
                    "subject": pdf_reader.metadata.get("/Subject", ""),
                    "creator": pdf_reader.metadata.get("/Creator", ""),
                    "producer": pdf_reader.metadata.get("/Producer", ""),
                    "creation_date": str(pdf_reader.metadata.get("/CreationDate", "")),
                    "modification_date": str(pdf_reader.metadata.get("/ModDate", ""))
                })
            
            # Extract text from all pages
            for page_num, page in enumerate(pdf_reader.pages):
                try:
                    page_text = page.extract_text()
                    text_content += f"\n--- Page {page_num + 1} ---\n{page_text}"
                except Exception as e:
                    logger.warning(f"Error extracting text from page {page_num + 1}: {e}")
                    continue
            
            return {
                "text_content": text_content.strip(),
                "metadata": metadata
            }
            
        except Exception as e:
            return {"error": f"PDF parsing error: {str(e)}"}
    
    async def _parse_docx(self, content: bytes) -> Dict[str, Any]:
        """Parse DOCX document"""
        try:
            import io
            doc = Document(io.BytesIO(content))
            
            text_content = ""
            metadata = {
                "num_paragraphs": len(doc.paragraphs),
                "num_tables": len(doc.tables),
                "num_sections": len(doc.sections)
            }
            
            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                text_content += paragraph.text + "\n"
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text_content += cell.text + " "
                    text_content += "\n"
            
            return {
                "text_content": text_content.strip(),
                "metadata": metadata
            }
            
        except Exception as e:
            return {"error": f"DOCX parsing error: {str(e)}"}
    
    async def _parse_excel(self, content: bytes) -> Dict[str, Any]:
        """Parse Excel document"""
        try:
            import io
            workbook = openpyxl.load_workbook(io.BytesIO(content))
            
            text_content = ""
            metadata = {
                "num_sheets": len(workbook.sheetnames),
                "sheet_names": workbook.sheetnames
            }
            
            # Extract text from all sheets
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_content += f"\n--- Sheet: {sheet_name} ---\n"
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = " ".join([str(cell) for cell in row if cell is not None])
                    if row_text.strip():
                        text_content += row_text + "\n"
            
            return {
                "text_content": text_content.strip(),
                "metadata": metadata
            }
            
        except Exception as e:
            return {"error": f"Excel parsing error: {str(e)}"}
    
    async def _parse_pptx(self, content: bytes) -> Dict[str, Any]:
        """Parse PowerPoint document"""
        try:
            import io
            prs = Presentation(io.BytesIO(content))
            
            text_content = ""
            metadata = {
                "num_slides": len(prs.slides)
            }
            
            # Extract text from all slides
            for slide_num, slide in enumerate(prs.slides):
                text_content += f"\n--- Slide {slide_num + 1} ---\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_content += shape.text + "\n"
            
            return {
                "text_content": text_content.strip(),
                "metadata": metadata
            }
            
        except Exception as e:
            return {"error": f"PowerPoint parsing error: {str(e)}"}
    
    async def _parse_txt(self, content: bytes) -> Dict[str, Any]:
        """Parse plain text document"""
        try:
            # Try different encodings
            encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
            
            for encoding in encodings:
                try:
                    text_content = content.decode(encoding)
                    break
                except UnicodeDecodeError:
                    continue
            else:
                # If all encodings fail, use utf-8 with error handling
                text_content = content.decode('utf-8', errors='replace')
            
            return {
                "text_content": text_content,
                "metadata": {"encoding": "detected"}
            }
            
        except Exception as e:
            return {"error": f"Text parsing error: {str(e)}"}
    
    async def _parse_image(self, content: bytes) -> Dict[str, Any]:
        """Parse image using OCR"""
        try:
            import io
            image = Image.open(io.BytesIO(content))
            
            # Convert to RGB if necessary
            if image.mode != 'RGB':
                image = image.convert('RGB')
            
            # Extract text using OCR
            text_content = pytesseract.image_to_string(image)
            
            metadata = {
                "image_size": image.size,
                "image_mode": image.mode,
                "ocr_confidence": "unknown"
            }
            
            return {
                "text_content": text_content.strip(),
                "metadata": metadata
            }
            
        except Exception as e:
            return {"error": f"OCR parsing error: {str(e)}"}
    
    def _is_image_file(self, filename: str) -> bool:
        """Check if file is an image"""
        image_extensions = ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp']
        return any(filename.lower().endswith(ext) for ext in image_extensions)
    
    def _generate_document_id(self, content: bytes, filename: str) -> str:
        """Generate unique document ID"""
        content_hash = hashlib.md5(content).hexdigest()
        filename_hash = hashlib.md5(filename.encode()).hexdigest()
        return f"{content_hash[:8]}_{filename_hash[:8]}"
    
    def _update_stats(self, processing_time: float, success: bool):
        """Update processing statistics"""
        self.stats["total_uploads"] += 1
        
        if success:
            self.stats["successful_uploads"] += 1
        else:
            self.stats["failed_uploads"] += 1
        
        # Update average processing time
        total_uploads = self.stats["total_uploads"]
        current_avg = self.stats["average_processing_time"]
        new_avg = ((current_avg * (total_uploads - 1)) + processing_time) / total_uploads
        self.stats["average_processing_time"] = round(new_avg, 3)
    
    def get_supported_formats(self) -> Dict[str, bool]:
        """Get list of supported file formats"""
        return self.supported_formats.copy()
    
    def is_format_supported(self, filename: str) -> bool:
        """Check if file format is supported"""
        file_extension = Path(filename).suffix.lower()
        
        if file_extension == '.pdf':
            return PDF_AVAILABLE
        elif file_extension == '.docx':
            return DOCX_AVAILABLE
        elif file_extension in ['.xlsx', '.xls']:
            return EXCEL_AVAILABLE
        elif file_extension == '.pptx':
            return PPTX_AVAILABLE
        elif file_extension == '.txt':
            return True
        elif self._is_image_file(filename):
            return OCR_AVAILABLE
        else:
            return False
    
    def get_stats(self) -> Dict[str, Any]:
        """Get processing statistics"""
        uptime = time.time() - self.stats["start_time"]
        return {
            "stats": self.stats.copy(),
            "uptime_seconds": round(uptime, 2),
            "uptime_hours": round(uptime / 3600, 2),
            "success_rate": round(self.stats["successful_uploads"] / max(1, self.stats["total_uploads"]) * 100, 2),
            "supported_formats": self.supported_formats
        }

# Global instance
document_upload_processor = DocumentUploadProcessor()













