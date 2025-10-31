"""
Real Document Parser
Handles various document formats with real parsing capabilities
"""

import logging
import io
from typing import Dict, List, Any, Optional, Union
from pathlib import Path
import mimetypes

# Document parsing libraries
try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from PIL import Image
    import pytesseract
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

logger = logging.getLogger(__name__)

class RealDocumentParser:
    """Real document parser for various formats"""
    
    def __init__(self):
        self.supported_formats = {
            'pdf': PDF_AVAILABLE,
            'docx': DOCX_AVAILABLE,
            'xlsx': EXCEL_AVAILABLE,
            'pptx': PPTX_AVAILABLE,
            'txt': True,
            'image': OCR_AVAILABLE
        }
    
    async def parse_document(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """Parse document and extract text content"""
        try:
            file_extension = Path(filename).suffix.lower()
            mime_type = mimetypes.guess_type(filename)[0]
            
            logger.info(f"Parsing document: {filename} (type: {mime_type})")
            
            result = {
                "filename": filename,
                "file_type": mime_type,
                "file_extension": file_extension,
                "file_size": len(file_content),
                "parsing_successful": False,
                "text_content": "",
                "metadata": {},
                "error": None
            }
            
            # Parse based on file type
            if file_extension == '.pdf' and PDF_AVAILABLE:
                result.update(await self._parse_pdf(file_content))
            elif file_extension == '.docx' and DOCX_AVAILABLE:
                result.update(await self._parse_docx(file_content))
            elif file_extension in ['.xlsx', '.xls'] and EXCEL_AVAILABLE:
                result.update(await self._parse_excel(file_content))
            elif file_extension == '.pptx' and PPTX_AVAILABLE:
                result.update(await self._parse_pptx(file_content))
            elif file_extension == '.txt':
                result.update(await self._parse_txt(file_content))
            elif mime_type and mime_type.startswith('image/') and OCR_AVAILABLE:
                result.update(await self._parse_image(file_content))
            else:
                result["error"] = f"Unsupported file format: {file_extension}"
                return result
            
            result["parsing_successful"] = True
            return result
            
        except Exception as e:
            logger.error(f"Error parsing document {filename}: {e}")
            return {
                "filename": filename,
                "parsing_successful": False,
                "error": str(e),
                "text_content": "",
                "metadata": {}
            }
    
    async def _parse_pdf(self, content: bytes) -> Dict[str, Any]:
        """Parse PDF document"""
        try:
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
            
            text_content = ""
            metadata = {
                "num_pages": len(pdf_reader.pages),
                "title": "",
                "author": "",
                "subject": "",
                "creator": "",
                "producer": "",
                "creation_date": "",
                "modification_date": ""
            }
            
            # Extract metadata
            if pdf_reader.metadata:
                metadata.update({
                    "title": pdf_reader.metadata.get("/Title", ""),
                    "author": pdf_reader.metadata.get("/Author", ""),
                    "subject": pdf_reader.metadata.get("/Subject", ""),
                    "creator": pdf_reader.metadata.get("/Creator", ""),
                    "producer": pdf_reader.metadata.get("/Producer", ""),
                    "creation_date": str(pdf_reader.metadata.get("/CreationDate", "")),
                    "modification_date": str(pdf_reader.metadata.get("/ModDate", ""))
                })
            
            # Extract text from all pages
            for page_num, page in enumerate(pdf_reader.pages):
                try:
                    page_text = page.extract_text()
                    text_content += f"\n--- Page {page_num + 1} ---\n{page_text}"
                except Exception as e:
                    logger.warning(f"Error extracting text from page {page_num + 1}: {e}")
                    continue
            
            return {
                "text_content": text_content.strip(),
                "metadata": metadata
            }
            
        except Exception as e:
            return {"error": f"PDF parsing error: {str(e)}"}
    
    async def _parse_docx(self, content: bytes) -> Dict[str, Any]:
        """Parse DOCX document"""
        try:
            doc = Document(io.BytesIO(content))
            
            text_content = ""
            metadata = {
                "num_paragraphs": len(doc.paragraphs),
                "num_tables": len(doc.tables),
                "num_sections": len(doc.sections)
            }
            
            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                text_content += paragraph.text + "\n"
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text_content += cell.text + " "
                    text_content += "\n"
            
            return {
                "text_content": text_content.strip(),
                "metadata": metadata
            }
            
        except Exception as e:
            return {"error": f"DOCX parsing error: {str(e)}"}
    
    async def _parse_excel(self, content: bytes) -> Dict[str, Any]:
        """Parse Excel document"""
        try:
            workbook = openpyxl.load_workbook(io.BytesIO(content))
            
            text_content = ""
            metadata = {
                "num_sheets": len(workbook.sheetnames),
                "sheet_names": workbook.sheetnames
            }
            
            # Extract text from all sheets
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_content += f"\n--- Sheet: {sheet_name} ---\n"
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = " ".join([str(cell) for cell in row if cell is not None])
                    if row_text.strip():
                        text_content += row_text + "\n"
            
            return {
                "text_content": text_content.strip(),
                "metadata": metadata
            }
            
        except Exception as e:
            return {"error": f"Excel parsing error: {str(e)}"}
    
    async def _parse_pptx(self, content: bytes) -> Dict[str, Any]:
        """Parse PowerPoint document"""
        try:
            prs = Presentation(io.BytesIO(content))
            
            text_content = ""
            metadata = {
                "num_slides": len(prs.slides)
            }
            
            # Extract text from all slides
            for slide_num, slide in enumerate(prs.slides):
                text_content += f"\n--- Slide {slide_num + 1} ---\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_content += shape.text + "\n"
            
            return {
                "text_content": text_content.strip(),
                "metadata": metadata
            }
            
        except Exception as e:
            return {"error": f"PowerPoint parsing error: {str(e)}"}
    
    async def _parse_txt(self, content: bytes) -> Dict[str, Any]:
        """Parse plain text document"""
        try:
            # Try different encodings
            encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
            
            for encoding in encodings:
                try:
                    text_content = content.decode(encoding)
                    break
                except UnicodeDecodeError:
                    continue
            else:
                # If all encodings fail, use utf-8 with error handling
                text_content = content.decode('utf-8', errors='replace')
            
            return {
                "text_content": text_content,
                "metadata": {"encoding": "detected"}
            }
            
        except Exception as e:
            return {"error": f"Text parsing error: {str(e)}"}
    
    async def _parse_image(self, content: bytes) -> Dict[str, Any]:
        """Parse image using OCR"""
        try:
            image = Image.open(io.BytesIO(content))
            
            # Convert to RGB if necessary
            if image.mode != 'RGB':
                image = image.convert('RGB')
            
            # Extract text using OCR
            text_content = pytesseract.image_to_string(image)
            
            metadata = {
                "image_size": image.size,
                "image_mode": image.mode,
                "ocr_confidence": "unknown"
            }
            
            return {
                "text_content": text_content.strip(),
                "metadata": metadata
            }
            
        except Exception as e:
            return {"error": f"OCR parsing error: {str(e)}"}
    
    def get_supported_formats(self) -> Dict[str, bool]:
        """Get list of supported file formats"""
        return self.supported_formats.copy()
    
    def is_format_supported(self, filename: str) -> bool:
        """Check if file format is supported"""
        file_extension = Path(filename).suffix.lower()
        
        if file_extension == '.pdf':
            return PDF_AVAILABLE
        elif file_extension == '.docx':
            return DOCX_AVAILABLE
        elif file_extension in ['.xlsx', '.xls']:
            return EXCEL_AVAILABLE
        elif file_extension == '.pptx':
            return PPTX_AVAILABLE
        elif file_extension == '.txt':
            return True
        elif file_extension in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff']:
            return OCR_AVAILABLE
        else:
            return False

# Global instance
document_parser = RealDocumentParser()













