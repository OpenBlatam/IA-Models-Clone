"""
Advanced AI Document Processing Engine
"""

import asyncio
import hashlib
import logging
import os
import tempfile
import time
from typing import Dict, Any, List, Optional, Tuple, BinaryIO
from datetime import datetime
import uuid

import numpy as np
import pandas as pd
from PIL import Image
import cv2
import easyocr
import pytesseract
from transformers import pipeline, AutoTokenizer, AutoModel
from sentence_transformers import SentenceTransformer
import spacy
import nltk
from textblob import TextBlob
from langdetect import detect
import gensim
from gensim import corpora, models
import torch
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
from sklearn.cluster import KMeans
from sklearn.decomposition import LatentDirichletAllocation

# Document processing libraries
import PyPDF2
import pdfplumber
from docx import Document
import openpyxl
from pptx import Presentation
import magic

from config import settings
from models import (
    DocumentType, ProcessingStatus, AnalysisType,
    DocumentMetadata, OCRResult, ClassificationResult,
    EntityExtractionResult, SentimentAnalysisResult,
    TopicModelingResult, SummarizationResult,
    KeywordExtractionResult, ContentAnalysisResult,
    DocumentProcessingResult
)

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """Advanced AI Document Processing Engine"""
    
    def __init__(self):
        self.models = {}
        self.processors = {}
        self.ocr_reader = None
        self.nlp = None
        self.initialized = False
        self.vectorizer = None
        self.sentence_transformer = None
        
    async def initialize(self):
        """Initialize the document processing engine"""
        if self.initialized:
            return
            
        try:
            logger.info("Initializing AI Document Processor...")
            
            # Initialize spaCy model
            try:
                self.nlp = spacy.load(settings.ner_model)
            except OSError:
                logger.warning("spaCy model not found. Installing...")
                import subprocess
                subprocess.run(["python", "-m", "spacy", "download", settings.ner_model])
                self.nlp = spacy.load(settings.ner_model)
            
            # Initialize NLTK data
            try:
                nltk.data.find('tokenizers/punkt')
                nltk.data.find('corpora/stopwords')
                nltk.data.find('vader_lexicon')
            except LookupError:
                logger.info("Downloading NLTK data...")
                nltk.download('punkt', quiet=True)
                nltk.download('stopwords', quiet=True)
                nltk.download('vader_lexicon', quiet=True)
            
            # Initialize sentence transformer
            self.sentence_transformer = SentenceTransformer(settings.embedding_model)
            
            # Initialize classification pipeline
            self.models['classifier'] = pipeline(
                "text-classification",
                model=settings.classification_model,
                return_all_scores=True
            )
            
            # Initialize summarization pipeline
            self.models['summarizer'] = pipeline(
                "summarization",
                model=settings.summarization_model
            )
            
            # Initialize sentiment analysis
            self.models['sentiment'] = pipeline(
                "sentiment-analysis",
                model="cardiffnlp/twitter-roberta-base-sentiment-latest",
                return_all_scores=True
            )
            
            # Initialize OCR
            self.ocr_reader = easyocr.Reader(settings.ocr_languages)
            
            # Initialize TF-IDF vectorizer
            self.vectorizer = TfidfVectorizer(
                max_features=1000,
                stop_words='english',
                ngram_range=(1, 2)
            )
            
            # Initialize topic modeling
            self.models['lda'] = LatentDirichletAllocation(
                n_components=10,
                random_state=42,
                max_iter=10
            )
            
            self.initialized = True
            logger.info("AI Document Processor initialized successfully")
            
        except Exception as e:
            logger.error(f"Error initializing document processor: {e}")
            raise
    
    async def process_document(self, file_path: str, analysis_types: List[AnalysisType], 
                             metadata: Dict[str, Any] = None) -> DocumentProcessingResult:
        """Process a document with specified analysis types"""
        if not self.initialized:
            await self.initialize()
        
        document_id = str(uuid.uuid4())
        start_time = time.time()
        
        try:
            # Extract text from document
            text_content = await self._extract_text(file_path)
            
            # Get document metadata
            doc_metadata = await self._extract_metadata(file_path, text_content, metadata)
            doc_metadata.document_id = document_id
            
            # Initialize result
            result = DocumentProcessingResult(
                document_id=document_id,
                status=ProcessingStatus.PROCESSING,
                metadata=doc_metadata,
                processing_time=0.0,
                created_at=datetime.now()
            )
            
            # Perform requested analyses
            if AnalysisType.OCR in analysis_types:
                result.ocr_result = await self._perform_ocr(file_path)
            
            if AnalysisType.CLASSIFICATION in analysis_types:
                result.classification_result = await self._classify_document(text_content)
            
            if AnalysisType.ENTITY_EXTRACTION in analysis_types:
                result.entity_result = await self._extract_entities(text_content)
            
            if AnalysisType.SENTIMENT_ANALYSIS in analysis_types:
                result.sentiment_result = await self._analyze_sentiment(text_content)
            
            if AnalysisType.TOPIC_MODELING in analysis_types:
                result.topic_result = await self._model_topics(text_content)
            
            if AnalysisType.SUMMARIZATION in analysis_types:
                result.summarization_result = await self._summarize_document(text_content)
            
            if AnalysisType.KEYWORD_EXTRACTION in analysis_types:
                result.keyword_result = await self._extract_keywords(text_content)
            
            if AnalysisType.CONTENT_ANALYSIS in analysis_types:
                result.content_analysis_result = await self._analyze_content(text_content)
            
            # Update result
            result.status = ProcessingStatus.COMPLETED
            result.processing_time = time.time() - start_time
            
            return result
            
        except Exception as e:
            logger.error(f"Error processing document {document_id}: {e}")
            return DocumentProcessingResult(
                document_id=document_id,
                status=ProcessingStatus.FAILED,
                metadata=DocumentMetadata(
                    document_id=document_id,
                    filename=os.path.basename(file_path),
                    content_type="unknown",
                    file_size=0,
                    language="unknown",
                    created_at=datetime.now()
                ),
                processing_time=time.time() - start_time,
                created_at=datetime.now(),
                error_message=str(e)
            )
    
    async def _extract_text(self, file_path: str) -> str:
        """Extract text from various document formats"""
        try:
            file_extension = os.path.splitext(file_path)[1].lower()
            
            if file_extension == '.pdf':
                return await self._extract_text_from_pdf(file_path)
            elif file_extension in ['.docx', '.doc']:
                return await self._extract_text_from_docx(file_path)
            elif file_extension == '.txt':
                return await self._extract_text_from_txt(file_path)
            elif file_extension in ['.xlsx', '.xls']:
                return await self._extract_text_from_excel(file_path)
            elif file_extension == '.pptx':
                return await self._extract_text_from_pptx(file_path)
            elif file_extension in ['.png', '.jpg', '.jpeg', '.tiff', '.bmp']:
                return await self._extract_text_from_image(file_path)
            else:
                # Try to read as text file
                with open(file_path, 'r', encoding='utf-8') as f:
                    return f.read()
                    
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {e}")
            return ""
    
    async def _extract_text_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF using pdfplumber"""
        try:
            text = ""
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
            return text
        except Exception as e:
            logger.error(f"Error extracting text from PDF: {e}")
            return ""
    
    async def _extract_text_from_docx(self, file_path: str) -> str:
        """Extract text from DOCX file"""
        try:
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except Exception as e:
            logger.error(f"Error extracting text from DOCX: {e}")
            return ""
    
    async def _extract_text_from_txt(self, file_path: str) -> str:
        """Extract text from TXT file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            logger.error(f"Error extracting text from TXT: {e}")
            return ""
    
    async def _extract_text_from_excel(self, file_path: str) -> str:
        """Extract text from Excel file"""
        try:
            workbook = openpyxl.load_workbook(file_path)
            text = ""
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value:
                            text += str(cell.value) + " "
                    text += "\n"
            return text
        except Exception as e:
            logger.error(f"Error extracting text from Excel: {e}")
            return ""
    
    async def _extract_text_from_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint file"""
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            logger.error(f"Error extracting text from PowerPoint: {e}")
            return ""
    
    async def _extract_text_from_image(self, file_path: str) -> str:
        """Extract text from image using OCR"""
        try:
            # Use EasyOCR for better accuracy
            results = self.ocr_reader.readtext(file_path)
            text = ""
            for (bbox, text_item, confidence) in results:
                if confidence > settings.ocr_confidence_threshold:
                    text += text_item + " "
            return text
        except Exception as e:
            logger.error(f"Error extracting text from image: {e}")
            return ""
    
    async def _extract_metadata(self, file_path: str, text_content: str, 
                              additional_metadata: Dict[str, Any] = None) -> DocumentMetadata:
        """Extract document metadata"""
        try:
            file_stats = os.stat(file_path)
            file_size = file_stats.st_size
            
            # Detect language
            language = "en"  # default
            if text_content:
                try:
                    language = detect(text_content[:1000])  # Use first 1000 chars
                except:
                    language = "en"
            
            # Count words and characters
            word_count = len(text_content.split()) if text_content else 0
            character_count = len(text_content) if text_content else 0
            
            # Get file info
            mime_type = magic.from_file(file_path, mime=True)
            filename = os.path.basename(file_path)
            
            return DocumentMetadata(
                document_id="",  # Will be set by caller
                filename=filename,
                content_type=mime_type,
                file_size=file_size,
                language=language,
                word_count=word_count,
                character_count=character_count,
                created_at=datetime.now(),
                modified_at=datetime.fromtimestamp(file_stats.st_mtime),
                **additional_metadata or {}
            )
            
        except Exception as e:
            logger.error(f"Error extracting metadata: {e}")
            return DocumentMetadata(
                document_id="",
                filename=os.path.basename(file_path),
                content_type="unknown",
                file_size=0,
                language="en",
                created_at=datetime.now()
            )
    
    async def _perform_ocr(self, file_path: str) -> OCRResult:
        """Perform OCR on document"""
        start_time = time.time()
        
        try:
            if file_path.lower().endswith('.pdf'):
                # Convert PDF to images and OCR each page
                text = ""
                confidence_scores = []
                
                # This is a simplified version - in production you'd use pdf2image
                with pdfplumber.open(file_path) as pdf:
                    for page_num, page in enumerate(pdf.pages):
                        # Extract text directly from PDF
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text
                            confidence_scores.append(0.9)  # PDF text extraction confidence
                
                avg_confidence = np.mean(confidence_scores) if confidence_scores else 0.0
                
            else:
                # OCR for images
                results = self.ocr_reader.readtext(file_path)
                text = ""
                confidence_scores = []
                
                for (bbox, text_item, confidence) in results:
                    if confidence > settings.ocr_confidence_threshold:
                        text += text_item + " "
                        confidence_scores.append(confidence)
                
                avg_confidence = np.mean(confidence_scores) if confidence_scores else 0.0
            
            # Detect language
            language = "en"
            if text:
                try:
                    language = detect(text[:1000])
                except:
                    language = "en"
            
            return OCRResult(
                text=text,
                confidence=avg_confidence,
                language=language,
                page_results=[],  # Would be populated in full implementation
                processing_time=time.time() - start_time
            )
            
        except Exception as e:
            logger.error(f"Error performing OCR: {e}")
            return OCRResult(
                text="",
                confidence=0.0,
                language="en",
                page_results=[],
                processing_time=time.time() - start_time
            )
    
    async def _classify_document(self, text: str) -> ClassificationResult:
        """Classify document content"""
        start_time = time.time()
        
        try:
            if not text.strip():
                return ClassificationResult(
                    category="unknown",
                    confidence=0.0,
                    all_categories=[],
                    processing_time=time.time() - start_time
                )
            
            # Truncate text if too long
            if len(text) > 512:
                text = text[:512]
            
            results = self.models['classifier'](text)
            
            # Get best category
            best_result = max(results[0], key=lambda x: x['score'])
            
            return ClassificationResult(
                category=best_result['label'],
                confidence=best_result['score'],
                all_categories=results[0],
                processing_time=time.time() - start_time
            )
            
        except Exception as e:
            logger.error(f"Error classifying document: {e}")
            return ClassificationResult(
                category="unknown",
                confidence=0.0,
                all_categories=[],
                processing_time=time.time() - start_time
            )
    
    async def _extract_entities(self, text: str) -> EntityExtractionResult:
        """Extract named entities from text"""
        start_time = time.time()
        
        try:
            if not text.strip():
                return EntityExtractionResult(
                    entities=[],
                    entity_count=0,
                    entity_types=[],
                    processing_time=time.time() - start_time
                )
            
            doc = self.nlp(text)
            entities = []
            entity_types = set()
            
            for ent in doc.ents:
                entities.append({
                    "text": ent.text,
                    "label": ent.label_,
                    "start": ent.start_char,
                    "end": ent.end_char,
                    "description": spacy.explain(ent.label_)
                })
                entity_types.add(ent.label_)
            
            return EntityExtractionResult(
                entities=entities,
                entity_count=len(entities),
                entity_types=list(entity_types),
                processing_time=time.time() - start_time
            )
            
        except Exception as e:
            logger.error(f"Error extracting entities: {e}")
            return EntityExtractionResult(
                entities=[],
                entity_count=0,
                entity_types=[],
                processing_time=time.time() - start_time
            )
    
    async def _analyze_sentiment(self, text: str) -> SentimentAnalysisResult:
        """Analyze sentiment of text"""
        start_time = time.time()
        
        try:
            if not text.strip():
                return SentimentAnalysisResult(
                    sentiment="neutral",
                    confidence=0.0,
                    polarity=0.0,
                    subjectivity=0.0,
                    processing_time=time.time() - start_time
                )
            
            # Use transformer-based sentiment analysis
            results = self.models['sentiment'](text)
            
            # Extract sentiment scores
            sentiment_scores = {}
            for result in results[0]:
                sentiment_scores[result['label']] = result['score']
            
            # Determine dominant sentiment
            dominant_sentiment = max(sentiment_scores, key=sentiment_scores.get)
            
            # Also use TextBlob for comparison
            blob = TextBlob(text)
            polarity = blob.sentiment.polarity
            subjectivity = blob.sentiment.subjectivity
            
            return SentimentAnalysisResult(
                sentiment=dominant_sentiment,
                confidence=sentiment_scores[dominant_sentiment],
                polarity=polarity,
                subjectivity=subjectivity,
                processing_time=time.time() - start_time
            )
            
        except Exception as e:
            logger.error(f"Error analyzing sentiment: {e}")
            return SentimentAnalysisResult(
                sentiment="neutral",
                confidence=0.0,
                polarity=0.0,
                subjectivity=0.0,
                processing_time=time.time() - start_time
            )
    
    async def _model_topics(self, text: str) -> TopicModelingResult:
        """Perform topic modeling on text"""
        start_time = time.time()
        
        try:
            if not text.strip():
                return TopicModelingResult(
                    topics=[],
                    topic_count=0,
                    dominant_topic="",
                    processing_time=time.time() - start_time
                )
            
            # Preprocess text
            doc = self.nlp(text)
            tokens = [token.lemma_.lower() for token in doc 
                     if not token.is_stop and not token.is_punct and token.is_alpha]
            processed_text = ' '.join(tokens)
            
            if not processed_text:
                return TopicModelingResult(
                    topics=[],
                    topic_count=0,
                    dominant_topic="",
                    processing_time=time.time() - start_time
                )
            
            # Create TF-IDF matrix
            tfidf_matrix = self.vectorizer.fit_transform([processed_text])
            
            # Fit LDA model
            lda_model = LatentDirichletAllocation(
                n_components=5,
                random_state=42,
                max_iter=10
            )
            lda_model.fit(tfidf_matrix)
            
            # Extract topics
            feature_names = self.vectorizer.get_feature_names_out()
            topics = []
            
            for topic_idx, topic in enumerate(lda_model.components_):
                top_words_idx = topic.argsort()[-10:][::-1]
                top_words = [feature_names[i] for i in top_words_idx]
                topics.append({
                    "topic_id": topic_idx,
                    "words": top_words,
                    "weights": topic[top_words_idx].tolist()
                })
            
            # Find dominant topic
            dominant_topic = topics[0]["words"][0] if topics else ""
            
            return TopicModelingResult(
                topics=topics,
                topic_count=len(topics),
                dominant_topic=dominant_topic,
                processing_time=time.time() - start_time
            )
            
        except Exception as e:
            logger.error(f"Error modeling topics: {e}")
            return TopicModelingResult(
                topics=[],
                topic_count=0,
                dominant_topic="",
                processing_time=time.time() - start_time
            )
    
    async def _summarize_document(self, text: str) -> SummarizationResult:
        """Summarize document content"""
        start_time = time.time()
        
        try:
            if not text.strip():
                return SummarizationResult(
                    summary="",
                    original_length=0,
                    summary_length=0,
                    compression_ratio=0.0,
                    key_points=[],
                    processing_time=time.time() - start_time
                )
            
            # Truncate text if too long
            if len(text) > 1024:
                text = text[:1024]
            
            summary = self.models['summarizer'](text, max_length=150, min_length=30, do_sample=False)
            summary_text = summary[0]["summary_text"]
            
            # Extract key points (simplified)
            sentences = text.split('.')
            key_points = [s.strip() for s in sentences[:3] if s.strip()]
            
            return SummarizationResult(
                summary=summary_text,
                original_length=len(text),
                summary_length=len(summary_text),
                compression_ratio=len(summary_text) / len(text),
                key_points=key_points,
                processing_time=time.time() - start_time
            )
            
        except Exception as e:
            logger.error(f"Error summarizing document: {e}")
            return SummarizationResult(
                summary="",
                original_length=len(text) if text else 0,
                summary_length=0,
                compression_ratio=0.0,
                key_points=[],
                processing_time=time.time() - start_time
            )
    
    async def _extract_keywords(self, text: str) -> KeywordExtractionResult:
        """Extract keywords from text"""
        start_time = time.time()
        
        try:
            if not text.strip():
                return KeywordExtractionResult(
                    keywords=[],
                    keyword_count=0,
                    processing_time=time.time() - start_time
                )
            
            # Use TF-IDF for keyword extraction
            tfidf_matrix = self.vectorizer.fit_transform([text])
            feature_names = self.vectorizer.get_feature_names_out()
            
            # Get top keywords
            scores = tfidf_matrix.toarray()[0]
            keyword_indices = scores.argsort()[-20:][::-1]  # Top 20 keywords
            
            keywords = []
            for idx in keyword_indices:
                if scores[idx] > 0:
                    keywords.append({
                        "word": feature_names[idx],
                        "score": float(scores[idx])
                    })
            
            return KeywordExtractionResult(
                keywords=keywords,
                keyword_count=len(keywords),
                processing_time=time.time() - start_time
            )
            
        except Exception as e:
            logger.error(f"Error extracting keywords: {e}")
            return KeywordExtractionResult(
                keywords=[],
                keyword_count=0,
                processing_time=time.time() - start_time
            )
    
    async def _analyze_content(self, text: str) -> ContentAnalysisResult:
        """Analyze content quality and readability"""
        start_time = time.time()
        
        try:
            if not text.strip():
                return ContentAnalysisResult(
                    readability_score=0.0,
                    complexity_score=0.0,
                    quality_rating="poor",
                    suggestions=[],
                    processing_time=time.time() - start_time
                )
            
            blob = TextBlob(text)
            
            # Basic metrics
            sentences = blob.sentences
            words = blob.words
            
            # Calculate readability metrics
            avg_sentence_length = sum(len(sentence.words) for sentence in sentences) / len(sentences) if sentences else 0
            avg_word_length = sum(len(word) for word in words) / len(words) if words else 0
            
            # Flesch Reading Ease (simplified)
            flesch_score = 206.835 - (1.015 * avg_sentence_length) - (84.6 * avg_word_length)
            
            # Grade level estimation
            grade_level = (0.39 * avg_sentence_length) + (11.8 * avg_word_length) - 15.59
            
            # Quality rating
            if flesch_score >= 80:
                quality_rating = "excellent"
            elif flesch_score >= 60:
                quality_rating = "good"
            elif flesch_score >= 40:
                quality_rating = "fair"
            else:
                quality_rating = "poor"
            
            # Generate suggestions
            suggestions = []
            if avg_sentence_length > 20:
                suggestions.append("Consider shortening sentences for better readability")
            if avg_word_length > 6:
                suggestions.append("Use simpler words where possible")
            if len(sentences) < 5:
                suggestions.append("Add more content to provide better context")
            
            return ContentAnalysisResult(
                readability_score=flesch_score,
                complexity_score=grade_level,
                quality_rating=quality_rating,
                suggestions=suggestions,
                processing_time=time.time() - start_time
            )
            
        except Exception as e:
            logger.error(f"Error analyzing content: {e}")
            return ContentAnalysisResult(
                readability_score=0.0,
                complexity_score=0.0,
                quality_rating="unknown",
                suggestions=[],
                processing_time=time.time() - start_time
            )


# Global document processor instance
document_processor = DocumentProcessor()


async def initialize_document_processor():
    """Initialize the document processor"""
    await document_processor.initialize()














