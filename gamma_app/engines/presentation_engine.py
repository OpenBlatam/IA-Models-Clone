"""
Gamma App - Presentation Engine
Advanced AI-powered presentation generation with design automation
"""

import asyncio
import json
import logging
from datetime import datetime
from typing import Dict, List, Optional, Any, Tuple
from dataclasses import dataclass, asdict
from enum import Enum
import io
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.units import inch
import requests
from PIL import Image

logger = logging.getLogger(__name__)

class SlideLayout(Enum):
    """Available slide layouts"""
    TITLE = "title"
    CONTENT = "content"
    TWO_COLUMN = "two_column"
    IMAGE_LEFT = "image_left"
    IMAGE_RIGHT = "image_right"
    IMAGE_CENTER = "image_center"
    CHART = "chart"
    QUOTE = "quote"
    COMPARISON = "comparison"
    TIMELINE = "timeline"
    CONCLUSION = "conclusion"

class ChartType(Enum):
    """Available chart types"""
    BAR = "bar"
    LINE = "line"
    PIE = "pie"
    COLUMN = "column"
    AREA = "area"
    SCATTER = "scatter"

@dataclass
class SlideDesign:
    """Design configuration for a slide"""
    layout: SlideLayout
    background_color: str
    text_color: str
    accent_color: str
    font_family: str
    font_size: int
    bullet_style: str
    image_position: str
    chart_type: Optional[ChartType] = None

@dataclass
class PresentationTheme:
    """Complete presentation theme"""
    name: str
    colors: Dict[str, str]
    fonts: Dict[str, str]
    layouts: Dict[str, SlideDesign]
    background_style: str
    logo_position: str

class PresentationEngine:
    """
    Advanced presentation generation engine
    """
    
    def __init__(self, config: Optional[Dict] = None):
        """Initialize the presentation engine"""
        self.config = config or {}
        self.themes = {}
        self.templates = {}
        self.image_cache = {}
        
        # Load themes and templates
        self._load_themes()
        self._load_templates()
        
        logger.info("Presentation Engine initialized successfully")

    def _load_themes(self):
        """Load available presentation themes"""
        self.themes = {
            "modern": PresentationTheme(
                name="Modern",
                colors={
                    "primary": "#2563eb",
                    "secondary": "#64748b",
                    "accent": "#f59e0b",
                    "background": "#ffffff",
                    "text": "#1e293b",
                    "light_bg": "#f8fafc"
                },
                fonts={
                    "heading": "Inter",
                    "body": "Inter",
                    "monospace": "JetBrains Mono"
                },
                layouts={},
                background_style="gradient",
                logo_position="top_right"
            ),
            "corporate": PresentationTheme(
                name="Corporate",
                colors={
                    "primary": "#1e40af",
                    "secondary": "#374151",
                    "accent": "#dc2626",
                    "background": "#ffffff",
                    "text": "#111827",
                    "light_bg": "#f9fafb"
                },
                fonts={
                    "heading": "Arial",
                    "body": "Arial",
                    "monospace": "Courier New"
                },
                layouts={},
                background_style="solid",
                logo_position="top_left"
            ),
            "creative": PresentationTheme(
                name="Creative",
                colors={
                    "primary": "#7c3aed",
                    "secondary": "#ec4899",
                    "accent": "#06b6d4",
                    "background": "#ffffff",
                    "text": "#1f2937",
                    "light_bg": "#fef3c7"
                },
                fonts={
                    "heading": "Poppins",
                    "body": "Open Sans",
                    "monospace": "Fira Code"
                },
                layouts={},
                background_style="pattern",
                logo_position="bottom_right"
            )
        }

    def _load_templates(self):
        """Load presentation templates"""
        self.templates = {
            "business_pitch": {
                "slide_count": 10,
                "sections": [
                    {"name": "title", "layout": SlideLayout.TITLE},
                    {"name": "problem", "layout": SlideLayout.CONTENT},
                    {"name": "solution", "layout": SlideLayout.TWO_COLUMN},
                    {"name": "market", "layout": SlideLayout.CHART},
                    {"name": "business_model", "layout": SlideLayout.IMAGE_LEFT},
                    {"name": "competition", "layout": SlideLayout.COMPARISON},
                    {"name": "team", "layout": SlideLayout.IMAGE_CENTER},
                    {"name": "financials", "layout": SlideLayout.CHART},
                    {"name": "timeline", "layout": SlideLayout.TIMELINE},
                    {"name": "conclusion", "layout": SlideLayout.CONCLUSION}
                ]
            },
            "academic": {
                "slide_count": 15,
                "sections": [
                    {"name": "title", "layout": SlideLayout.TITLE},
                    {"name": "abstract", "layout": SlideLayout.CONTENT},
                    {"name": "introduction", "layout": SlideLayout.CONTENT},
                    {"name": "literature_review", "layout": SlideLayout.TWO_COLUMN},
                    {"name": "methodology", "layout": SlideLayout.IMAGE_LEFT},
                    {"name": "data_collection", "layout": SlideLayout.CHART},
                    {"name": "results_1", "layout": SlideLayout.CHART},
                    {"name": "results_2", "layout": SlideLayout.CHART},
                    {"name": "analysis", "layout": SlideLayout.TWO_COLUMN},
                    {"name": "discussion", "layout": SlideLayout.CONTENT},
                    {"name": "limitations", "layout": SlideLayout.CONTENT},
                    {"name": "future_work", "layout": SlideLayout.CONTENT},
                    {"name": "conclusion", "layout": SlideLayout.CONCLUSION},
                    {"name": "references", "layout": SlideLayout.CONTENT},
                    {"name": "questions", "layout": SlideLayout.QUOTE}
                ]
            },
            "training": {
                "slide_count": 12,
                "sections": [
                    {"name": "title", "layout": SlideLayout.TITLE},
                    {"name": "agenda", "layout": SlideLayout.CONTENT},
                    {"name": "objectives", "layout": SlideLayout.CONTENT},
                    {"name": "overview", "layout": SlideLayout.IMAGE_CENTER},
                    {"name": "module_1", "layout": SlideLayout.CONTENT},
                    {"name": "module_2", "layout": SlideLayout.CONTENT},
                    {"name": "module_3", "layout": SlideLayout.CONTENT},
                    {"name": "examples", "layout": SlideLayout.IMAGE_LEFT},
                    {"name": "exercises", "layout": SlideLayout.TWO_COLUMN},
                    {"name": "best_practices", "layout": SlideLayout.QUOTE},
                    {"name": "resources", "layout": SlideLayout.CONTENT},
                    {"name": "conclusion", "layout": SlideLayout.CONCLUSION}
                ]
            }
        }

    async def create_presentation(self, content: Dict[str, Any], 
                                theme: str = "modern",
                                template: str = "business_pitch") -> bytes:
        """Create a PowerPoint presentation from content"""
        try:
            # Create new presentation
            prs = Presentation()
            
            # Apply theme
            theme_config = self.themes.get(theme, self.themes["modern"])
            template_config = self.templates.get(template, self.templates["business_pitch"])
            
            # Create slides based on template
            for section in template_config["sections"]:
                slide = await self._create_slide(
                    prs, content, section, theme_config
                )
            
            # Save to bytes
            output = io.BytesIO()
            prs.save(output)
            output.seek(0)
            
            logger.info(f"Presentation created successfully with {len(prs.slides)} slides")
            return output.getvalue()
            
        except Exception as e:
            logger.error(f"Error creating presentation: {e}")
            raise

    async def _create_slide(self, prs: Presentation, content: Dict[str, Any], 
                          section: Dict[str, Any], theme: PresentationTheme) -> None:
        """Create a single slide"""
        # Add slide with appropriate layout
        slide_layout = self._get_slide_layout(prs, section["layout"])
        slide = prs.slides.add_slide(slide_layout)
        
        # Apply theme colors
        self._apply_theme_colors(slide, theme)
        
        # Add content based on section
        if section["name"] == "title":
            await self._add_title_slide(slide, content, theme)
        elif section["layout"] == SlideLayout.CONTENT:
            await self._add_content_slide(slide, content, section, theme)
        elif section["layout"] == SlideLayout.TWO_COLUMN:
            await self._add_two_column_slide(slide, content, section, theme)
        elif section["layout"] == SlideLayout.CHART:
            await self._add_chart_slide(slide, content, section, theme)
        elif section["layout"] == SlideLayout.IMAGE_LEFT:
            await self._add_image_slide(slide, content, section, theme, "left")
        elif section["layout"] == SlideLayout.IMAGE_RIGHT:
            await self._add_image_slide(slide, content, section, theme, "right")
        elif section["layout"] == SlideLayout.IMAGE_CENTER:
            await self._add_image_slide(slide, content, section, theme, "center")
        elif section["layout"] == SlideLayout.COMPARISON:
            await self._add_comparison_slide(slide, content, section, theme)
        elif section["layout"] == SlideLayout.TIMELINE:
            await self._add_timeline_slide(slide, content, section, theme)
        elif section["layout"] == SlideLayout.QUOTE:
            await self._add_quote_slide(slide, content, section, theme)
        elif section["layout"] == SlideLayout.CONCLUSION:
            await self._add_conclusion_slide(slide, content, section, theme)

    def _get_slide_layout(self, prs: Presentation, layout_type: SlideLayout):
        """Get appropriate slide layout"""
        layout_map = {
            SlideLayout.TITLE: 0,  # Title slide layout
            SlideLayout.CONTENT: 1,  # Title and content layout
            SlideLayout.TWO_COLUMN: 2,  # Two content layout
            SlideLayout.IMAGE_LEFT: 3,  # Content with caption layout
            SlideLayout.IMAGE_RIGHT: 3,  # Content with caption layout
            SlideLayout.IMAGE_CENTER: 4,  # Blank layout
            SlideLayout.CHART: 1,  # Title and content layout
            SlideLayout.QUOTE: 4,  # Blank layout
            SlideLayout.COMPARISON: 2,  # Two content layout
            SlideLayout.TIMELINE: 1,  # Title and content layout
            SlideLayout.CONCLUSION: 0  # Title slide layout
        }
        
        layout_index = layout_map.get(layout_type, 1)
        return prs.slide_layouts[layout_index]

    def _apply_theme_colors(self, slide, theme: PresentationTheme):
        """Apply theme colors to slide"""
        # This would apply colors to slide background and text
        # Implementation depends on specific requirements
        pass

    async def _add_title_slide(self, slide, content: Dict[str, Any], theme: PresentationTheme):
        """Add title slide content"""
        # Add title
        title_shape = slide.shapes.title
        title_shape.text = content.get('title', 'Presentation Title')
        
        # Add subtitle if available
        if 'subtitle' in content:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = content['subtitle']
        
        # Apply title formatting
        self._format_title_text(title_shape, theme)

    async def _add_content_slide(self, slide, content: Dict[str, Any], 
                               section: Dict[str, Any], theme: PresentationTheme):
        """Add content slide"""
        # Add title
        title_shape = slide.shapes.title
        title_shape.text = section["name"].replace("_", " ").title()
        
        # Add content
        content_shape = slide.placeholders[1]
        content_text = self._get_section_content(content, section["name"])
        content_shape.text = content_text
        
        # Format content
        self._format_content_text(content_shape, theme)

    async def _add_two_column_slide(self, slide, content: Dict[str, Any], 
                                  section: Dict[str, Any], theme: PresentationTheme):
        """Add two-column slide"""
        # Add title
        title_shape = slide.shapes.title
        title_shape.text = section["name"].replace("_", " ").title()
        
        # Add content to both columns
        content_text = self._get_section_content(content, section["name"])
        content_parts = self._split_content_for_columns(content_text)
        
        # Left column
        left_shape = slide.placeholders[1]
        left_shape.text = content_parts[0]
        
        # Right column
        right_shape = slide.placeholders[2]
        right_shape.text = content_parts[1]

    async def _add_chart_slide(self, slide, content: Dict[str, Any], 
                             section: Dict[str, Any], theme: PresentationTheme):
        """Add chart slide"""
        # Add title
        title_shape = slide.shapes.title
        title_shape.text = section["name"].replace("_", " ").title()
        
        # Add chart placeholder
        chart_shape = slide.placeholders[1]
        chart_shape.text = f"[Chart placeholder for {section['name']}]"
        
        # In a real implementation, this would create actual charts
        # using libraries like matplotlib or plotly

    async def _add_image_slide(self, slide, content: Dict[str, Any], 
                             section: Dict[str, Any], theme: PresentationTheme, 
                             position: str):
        """Add image slide"""
        # Add title
        title_shape = slide.shapes.title
        title_shape.text = section["name"].replace("_", " ").title()
        
        # Add image placeholder
        if position == "left":
            image_shape = slide.placeholders[1]
            content_shape = slide.placeholders[2]
        elif position == "right":
            content_shape = slide.placeholders[1]
            image_shape = slide.placeholders[2]
        else:  # center
            image_shape = slide.placeholders[1]
            content_shape = slide.placeholders[2]
        
        # Add image placeholder text
        image_shape.text = f"[Image placeholder for {section['name']}]"
        
        # Add content
        content_text = self._get_section_content(content, section["name"])
        content_shape.text = content_text

    async def _add_comparison_slide(self, slide, content: Dict[str, Any], 
                                  section: Dict[str, Any], theme: PresentationTheme):
        """Add comparison slide"""
        # Add title
        title_shape = slide.shapes.title
        title_shape.text = section["name"].replace("_", " ").title()
        
        # Add comparison content
        content_text = self._get_section_content(content, section["name"])
        comparison_parts = self._split_content_for_comparison(content_text)
        
        # Left side
        left_shape = slide.placeholders[1]
        left_shape.text = comparison_parts[0]
        
        # Right side
        right_shape = slide.placeholders[2]
        right_shape.text = comparison_parts[1]

    async def _add_timeline_slide(self, slide, content: Dict[str, Any], 
                                section: Dict[str, Any], theme: PresentationTheme):
        """Add timeline slide"""
        # Add title
        title_shape = slide.shapes.title
        title_shape.text = section["name"].replace("_", " ").title()
        
        # Add timeline content
        content_shape = slide.placeholders[1]
        timeline_content = self._format_timeline_content(content, section["name"])
        content_shape.text = timeline_content

    async def _add_quote_slide(self, slide, content: Dict[str, Any], 
                             section: Dict[str, Any], theme: PresentationTheme):
        """Add quote slide"""
        # Add quote
        quote_shape = slide.placeholders[0]
        quote_text = self._get_section_content(content, section["name"])
        quote_shape.text = f'"{quote_text}"'
        
        # Format as quote
        self._format_quote_text(quote_shape, theme)

    async def _add_conclusion_slide(self, slide, content: Dict[str, Any], 
                                  section: Dict[str, Any], theme: PresentationTheme):
        """Add conclusion slide"""
        # Add title
        title_shape = slide.shapes.title
        title_shape.text = "Conclusion"
        
        # Add conclusion content
        content_shape = slide.placeholders[1]
        conclusion_content = self._get_section_content(content, section["name"])
        content_shape.text = conclusion_content
        
        # Add "Thank You" or "Questions?" at the bottom
        if len(slide.placeholders) > 2:
            footer_shape = slide.placeholders[2]
            footer_shape.text = "Thank You for Your Attention"

    def _get_section_content(self, content: Dict[str, Any], section_name: str) -> str:
        """Get content for a specific section"""
        if 'slides' in content:
            for slide in content['slides']:
                if slide.get('title', '').lower().replace(' ', '_') == section_name:
                    return slide.get('content', '')
        
        # Fallback to generic content
        return f"Content for {section_name.replace('_', ' ').title()}"

    def _split_content_for_columns(self, content: str) -> Tuple[str, str]:
        """Split content for two-column layout"""
        sentences = content.split('. ')
        mid_point = len(sentences) // 2
        left_content = '. '.join(sentences[:mid_point])
        right_content = '. '.join(sentences[mid_point:])
        return left_content, right_content

    def _split_content_for_comparison(self, content: str) -> Tuple[str, str]:
        """Split content for comparison layout"""
        # This would implement logic to split content into comparison parts
        # For now, simple split
        return self._split_content_for_columns(content)

    def _format_timeline_content(self, content: Dict[str, Any], section_name: str) -> str:
        """Format content as timeline"""
        timeline_content = self._get_section_content(content, section_name)
        # Add timeline formatting
        return f"• {timeline_content}"

    def _format_title_text(self, shape, theme: PresentationTheme):
        """Format title text"""
        # Apply title formatting based on theme
        pass

    def _format_content_text(self, shape, theme: PresentationTheme):
        """Format content text"""
        # Apply content formatting based on theme
        pass

    def _format_quote_text(self, shape, theme: PresentationTheme):
        """Format quote text"""
        # Apply quote formatting based on theme
        pass

    async def create_pdf_presentation(self, content: Dict[str, Any], 
                                    theme: str = "modern") -> bytes:
        """Create PDF version of presentation"""
        try:
            output = io.BytesIO()
            doc = SimpleDocTemplate(output, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            # Add title page
            title_style = styles['Title']
            title_style.fontSize = 24
            title_style.alignment = 1  # Center alignment
            
            story.append(Paragraph(content.get('title', 'Presentation'), title_style))
            story.append(Spacer(1, 12))
            
            if 'subtitle' in content:
                subtitle_style = styles['Heading2']
                subtitle_style.alignment = 1
                story.append(Paragraph(content['subtitle'], subtitle_style))
                story.append(Spacer(1, 12))
            
            story.append(PageBreak())
            
            # Add slides as pages
            if 'slides' in content:
                for slide in content['slides']:
                    # Slide title
                    slide_title = slide.get('title', 'Slide')
                    story.append(Paragraph(slide_title, styles['Heading1']))
                    story.append(Spacer(1, 12))
                    
                    # Slide content
                    slide_content = slide.get('content', '')
                    story.append(Paragraph(slide_content, styles['Normal']))
                    story.append(PageBreak())
            
            doc.build(story)
            output.seek(0)
            
            logger.info("PDF presentation created successfully")
            return output.getvalue()
            
        except Exception as e:
            logger.error(f"Error creating PDF presentation: {e}")
            raise

    async def add_images_to_presentation(self, presentation_bytes: bytes, 
                                       images: List[Dict[str, Any]]) -> bytes:
        """Add images to existing presentation"""
        # This would implement image addition to presentations
        # Implementation would depend on specific requirements
        return presentation_bytes

    async def add_charts_to_presentation(self, presentation_bytes: bytes, 
                                       chart_data: List[Dict[str, Any]]) -> bytes:
        """Add charts to existing presentation"""
        # This would implement chart addition to presentations
        # Implementation would depend on specific requirements
        return presentation_bytes

    def get_available_themes(self) -> List[str]:
        """Get available presentation themes"""
        return list(self.themes.keys())

    def get_available_templates(self) -> List[str]:
        """Get available presentation templates"""
        return list(self.templates.keys())

    def get_available_layouts(self) -> List[SlideLayout]:
        """Get available slide layouts"""
        return list(SlideLayout)

    def get_available_chart_types(self) -> List[ChartType]:
        """Get available chart types"""
        return list(ChartType)

    async def preview_slide(self, content: Dict[str, Any], slide_index: int, 
                          theme: str = "modern") -> bytes:
        """Generate preview of a specific slide"""
        # This would generate a preview image of a specific slide
        # Implementation would depend on specific requirements
        pass

    async def export_slide_as_image(self, presentation_bytes: bytes, 
                                  slide_index: int) -> bytes:
        """Export a specific slide as an image"""
        # This would export a slide as PNG/JPG
        # Implementation would depend on specific requirements
        pass



























